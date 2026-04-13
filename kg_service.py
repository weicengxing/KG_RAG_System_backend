"""
知识图谱RAG系统核心服务
包含文档处理、图谱构建、向量存储、RAG问答等功能
"""

import os
import json
import asyncio
import logging
from typing import List, Dict, Any, Tuple, Optional
from concurrent.futures import ThreadPoolExecutor
import time
import hashlib
from collections import defaultdict


# 文档处理
from PyPDF2 import PdfReader

# 文本分块器 - 兼容不同版本的langchain
try:
    # 新版本 langchain
    from langchain_text_splitters import RecursiveCharacterTextSplitter
except ImportError:
    try:
        # 旧版本 langchain
        from langchain.text_splitter import RecursiveCharacterTextSplitter
    except ImportError:
        # 如果都不行，使用简单的分块实现
        class RecursiveCharacterTextSplitter:
            def __init__(self, chunk_size, chunk_overlap, **kwargs):
                self.chunk_size = chunk_size
                self.chunk_overlap = chunk_overlap

            def split_text(self, text):
                """简单的文本分块实现"""
                chunks = []
                start = 0
                text_length = len(text)

                while start < text_length:
                    end = start + self.chunk_size
                    chunk = text[start:end]
                    chunks.append(chunk)
                    start = end - self.chunk_overlap

                return chunks

# 向量数据库 - ChromaDB (备用)
import chromadb
from chromadb.config import Settings

# 向量数据库 - Milvus (主用，高QPS)
try:
    from milvus_utils import MilvusVectorStore, get_milvus_client, MILVUS_ENABLED
    MILVUS_AVAILABLE = True
except ImportError:
    MILVUS_AVAILABLE = False
    print("[MILVUS] milvus_utils 未安装，跳过 Milvus 功能")

# Neo4j图数据库
from neo4j import GraphDatabase

# LLM和Embeddings
from openai import OpenAI
import requests

# Reranker
try:
    from FlagEmbedding import FlagReranker
    RERANKER_AVAILABLE = True
except ImportError:
    RERANKER_AVAILABLE = False
    print("[RERANKER] FlagEmbedding 未安装，跳过 Reranking 功能")

# Redis 客户端（用于缓存）
try:
    import redis
    REDIS_AVAILABLE = True
except ImportError:
    REDIS_AVAILABLE = False
    print("[CACHE] redis 库未安装，跳过缓存功能")


# 引入配置变量
DISABLE_VECTOR_SEARCH = os.getenv("DISABLE_VECTOR_SEARCH", "").lower() in ("1", "true", "yes", "on")
from config import (
    NEO4J_URI, NEO4J_USERNAME, NEO4J_PASSWORD,
    EXTRACTION_API_KEY, EXTRACTION_BASE_URL, EXTRACTION_MODEL,  # 三元组提取相关
    QA_MODELS, DEFAULT_QA_MODEL,  # AI问答模型列表
    EMBED_API_KEY, EMBED_BASE_URL, EMBED_MODEL,  # Embedding 相关
    CHUNK_SIZE, CHUNK_OVERLAP,
    VECTOR_SEARCH_TOP_K, VECTOR_SIMILARITY_THRESHOLD, GRAPH_SEARCH_HOPS,
    VECTOR_SEARCH_ENABLED, VECTOR_SEARCH_IN_THREAD,
    RERANKER_ENABLED, RERANKER_MODEL, RERANKER_TOP_K, RERANKER_INITIAL_TOP_K,
    RERANKER_SCORE_THRESHOLD, RERANKER_MAX_LENGTH, RERANKER_WEIGHT, VECTOR_WEIGHT,
    RAG_CACHE_ENABLED, RAG_CACHE_TTL,  # 缓存相关配置
    REDIS_HOST, REDIS_PORT, REDIS_DB, REDIS_PASSWORD,  # Redis连接配置
    QUERY_REWRITE_ENABLED, QUERY_REWRITE_COUNT, QUERY_REWRITE_TOP_K, QUERY_REWRITE_MERGE_TOP_K,  # Query Rewrite 配置
    BM25_ENABLED, BM25_TOP_K, BM25_WEIGHT, VECTOR_WEIGHT_HYBRID, RRF_K,  # BM25 混合检索配置
    MILVUS_ENABLED  # Milvus 向量数据库开关
)

# 导入 BM25 搜索工具
try:
    from elasticsearch_utils import index_kg_documents, search_kg_documents, delete_kg_documents
    ES_AVAILABLE = True
except ImportError:
    ES_AVAILABLE = False
    print("[BM25] elasticsearch_utils 未安装，跳过 BM25 检索功能")


# ==================== RAG 检索缓存 ====================

RAG_CACHE_PREFIX = "rag_cache:"  # 缓存键前缀


def get_redis_client():
    """获取 Redis 客户端（用于缓存）"""
    if not REDIS_AVAILABLE:
        return None
    try:
        client = redis.Redis(
            host=REDIS_HOST,
            port=REDIS_PORT,
            db=REDIS_DB,
            password=REDIS_PASSWORD,
            decode_responses=True
        )
        client.ping()
        return client
    except Exception as e:
        print(f"[CACHE] Redis 连接失败: {e}")
        return None


def cosine_similarity(vec1: List[float], vec2: List[float]) -> float:
    """计算两个向量的余弦相似度"""
    import numpy as np
    vec1 = np.array(vec1)
    vec2 = np.array(vec2)
    return float(np.dot(vec1, vec2) / (np.linalg.norm(vec1) * np.linalg.norm(vec2)))


def normalize_query(query: str) -> str:
    """标准化 Query：转小写、去标点、去空格"""
    # 1. 转小写
    query = query.lower()
    # 2. 去除标点符号（保留中文、英文、数字、下划线）
    import re
    query = re.sub(r'[^\w\s\u4e00-\u9fff]', '', query)
    # 3. 去除多余空格
    query = ' '.join(query.split())
    return query


def get_cache_key(query: str) -> str:
    """生成缓存键（基于标准化的查询内容）"""
    # 标准化 Query
    normalized_query = normalize_query(query)
    query_hash = hashlib.md5(normalized_query.encode('utf-8')).hexdigest()
    return f"{RAG_CACHE_PREFIX}{query_hash}"


def get_search_cache(query: str, query_embedding: Optional[List[float]] = None) -> Optional[List[Dict[str, Any]]]:
    """
    获取检索结果缓存（支持语义相似度匹配）

    Args:
        query: 查询文本
        query_embedding: 查询的 embedding 向量（可选，如果提供则启用语义缓存）

    Returns:
        缓存的检索结果，如果未命中则返回 None
    """
    if not RAG_CACHE_ENABLED or not REDIS_AVAILABLE:
        return None

    redis_client = get_redis_client()
    if not redis_client:
        return None

    try:
        # 1. 精确匹配缓存
        cache_key = get_cache_key(query)
        cached_data = redis_client.get(cache_key)
        if cached_data:
            print(f"[CACHE] 精确命中: {query[:30]}...")
            return json.loads(cached_data)

        # 2. 语义相似度缓存（如果启用且提供了 embedding）
        if SEMANTIC_CACHE_ENABLED and query_embedding:
            # 获取所有缓存的 query embeddings
            embedding_keys = redis_client.keys(f"{RAG_CACHE_PREFIX}emb:*")

            best_similarity = 0.0
            best_match_key = None

            for emb_key in embedding_keys[:SEMANTIC_CACHE_MAX_ENTRIES]:  # 限制检查数量
                try:
                    cached_emb_data = redis_client.get(emb_key)
                    if cached_emb_data:
                        cached_emb = json.loads(cached_emb_data)
                        similarity = cosine_similarity(query_embedding, cached_emb["embedding"])

                        if similarity > best_similarity:
                            best_similarity = similarity
                            best_match_key = cached_emb["result_key"]
                except Exception as e:
                    continue

            # 如果找到相似度超过阈值的缓存
            if best_similarity >= SEMANTIC_CACHE_SIMILARITY_THRESHOLD and best_match_key:
                cached_result = redis_client.get(best_match_key)
                if cached_result:
                    print(f"[CACHE] 语义命中: {query[:30]}... (相似度: {best_similarity:.3f})")
                    return json.loads(cached_result)

    except Exception as e:
        print(f"[CACHE] 读取缓存失败: {e}")

    return None


def set_search_cache(query: str, results: List[Dict[str, Any]], query_embedding: Optional[List[float]] = None) -> bool:
    """
    设置检索结果缓存（支持语义相似度缓存）

    Args:
        query: 查询文本
        results: 检索结果
        query_embedding: 查询的 embedding 向量（可选，如果提供则同时保存语义缓存）

    Returns:
        是否保存成功
    """
    if not RAG_CACHE_ENABLED or not REDIS_AVAILABLE:
        return False

    redis_client = get_redis_client()
    if not redis_client:
        return False

    try:
        # 1. 保存精确匹配缓存
        cache_key = get_cache_key(query)
        redis_client.setex(cache_key, RAG_CACHE_TTL, json.dumps(results, ensure_ascii=False))

        # 2. 保存语义缓存（如果启用且提供了 embedding）
        if SEMANTIC_CACHE_ENABLED and query_embedding:
            emb_key = f"{RAG_CACHE_PREFIX}emb:{hashlib.md5(query.encode('utf-8')).hexdigest()}"
            emb_data = {
                "query": query,
                "embedding": query_embedding,
                "result_key": cache_key
            }
            redis_client.setex(emb_key, RAG_CACHE_TTL, json.dumps(emb_data, ensure_ascii=False))
            print(f"[CACHE] 缓存已保存（含语义）: {query[:30]}... (TTL: {RAG_CACHE_TTL}s)")
        else:
            print(f"[CACHE] 缓存已保存: {query[:30]}... (TTL: {RAG_CACHE_TTL}s)")

        return True
    except Exception as e:
        print(f"[CACHE] 保存缓存失败: {e}")
        return False


# ==================== 知识图谱服务类 ====================

class KnowledgeGraphService:
    """知识图谱RAG服务类"""

    def __init__(self):
        """初始化服务"""
        # Neo4j连接
        self.neo4j_driver = GraphDatabase.driver(
            NEO4J_URI,
            auth=(NEO4J_USERNAME, NEO4J_PASSWORD)
        )

        # ChromaDB客户端（持久化存储）- 仅在未启用 Milvus 时使用
        self.chroma_client = None
        self.collection = None
        if not MILVUS_ENABLED or not MILVUS_AVAILABLE:
            chroma_data_dir = os.path.join(os.path.dirname(__file__), "chroma_data")
            os.makedirs(chroma_data_dir, exist_ok=True)
            self.chroma_client = chromadb.PersistentClient(
                path=chroma_data_dir,
                settings=Settings(anonymized_telemetry=False)
            )
            self.collection = self.chroma_client.get_or_create_collection(
                name="documents",
                metadata={"hnsw:space": "cosine"}
            )
            logging.info("[VECTOR] 使用 ChromaDB 向量存储")
        else:
            logging.info("[VECTOR] 使用 Milvus 向量存储")

        # Milvus 客户端（高QPS场景）
        self.milvus_client = None
        if MILVUS_ENABLED and MILVUS_AVAILABLE:
            try:
                from milvus_utils import get_milvus_client
                self.milvus_client = get_milvus_client()
                if self.milvus_client and self.milvus_client.connected:
                    logging.info("[MILVUS] ✅ Milvus 客户端初始化成功")
                else:
                    logging.warning("[MILVUS] Milvus 客户端连接失败，将回退到 ChromaDB")
                    # 回退到 ChromaDB
                    self._init_chromadb_fallback()
            except Exception as e:
                logging.warning(f"[MILVUS] 初始化失败: {e}，回退到 ChromaDB")
                self._init_chromadb_fallback()
        else:
            # 未启用 Milvus，使用 ChromaDB
            if not MILVUS_ENABLED:
                logging.info("[VECTOR] MILVUS_ENABLED=false，使用 ChromaDB")
            self._init_chromadb_fallback()

        # 初始化三元组提取客户端 (专用于实体关系抽取)
        self.extraction_client = OpenAI(
            api_key=EXTRACTION_API_KEY,
            base_url=EXTRACTION_BASE_URL
        )

        # 初始化 AI 问答客户端字典 (根据模型名称快速查找)
        self.qa_clients = {}
        for model_config in QA_MODELS:
            model_name = model_config["name"]
            self.qa_clients[model_name] = {
                "client": OpenAI(
                    api_key=model_config["api_key"],
                    base_url=model_config["base_url"]
                ),
                "model": model_config["model"],
                "use_qidian_api": model_config.get("use_qidian_api", False),
                "satoken": model_config.get("satoken", ""),
                "model_id": model_config.get("model_id", 32),
                "session_id": model_config.get("session_id", ""),
                "use_yansd_api": model_config.get("use_yansd_api", False),
                "session_cookie": model_config.get("session_cookie", ""),
                "new_api_user": model_config.get("new_api_user", "")
            }

        # 初始化 Embedding 客户端 (使用 ModelScope 配置)
        self.embed_client = OpenAI(
            api_key=EMBED_API_KEY,
            base_url=EMBED_BASE_URL
        )

        # 文本分块器
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=CHUNK_SIZE,
            chunk_overlap=CHUNK_OVERLAP,
            length_function=len,
            separators=["\n\n", "\n", "。", "！", "？", "；", ".", "!", "?", ";", " ", ""]
        )

        # 线程池用于并发处理
        self.executor = ThreadPoolExecutor(max_workers=10)
        self.pipeline_cache = {}
        self.pipeline_cache_ttl = 1800

        # 初始化 Reranker (延迟加载，首次使用时初始化)
        self.reranker = None
        if RERANKER_ENABLED and RERANKER_AVAILABLE:
            try:
                print(f"[RERANKER] 初始化 Reranker 模型: {RERANKER_MODEL}")
                self.reranker = FlagReranker(RERANKER_MODEL, use_fp16=True)
                print(f"[RERANKER] ✅ Reranker 初始化成功")
            except Exception as e:
                print(f"[RERANKER] ❌ Reranker 初始化失败: {e}")
                self.reranker = None

    def _init_chromadb_fallback(self):
        """回退到 ChromaDB（当 Milvus 不可用时）"""
        if self.chroma_client is None:
            chroma_data_dir = os.path.join(os.path.dirname(__file__), "chroma_data")
            os.makedirs(chroma_data_dir, exist_ok=True)
            self.chroma_client = chromadb.PersistentClient(
                path=chroma_data_dir,
                settings=Settings(anonymized_telemetry=False)
            )
        if self.collection is None:
            self.collection = self.chroma_client.get_or_create_collection(
                name="documents",
                metadata={"hnsw:space": "cosine"}
            )
        logging.info("[VECTOR] 已回退到 ChromaDB 向量存储")

    def __del__(self):
        """清理资源"""
        if hasattr(self, 'neo4j_driver'):
            self.neo4j_driver.close()

    def _get_pipeline_cache(self, doc_id: Optional[str]) -> Optional[Dict[str, Any]]:
        """获取文档处理缓存，超时则自动清理"""
        if not doc_id:
            return None

        cached = self.pipeline_cache.get(doc_id)
        if not cached:
            return None

        if time.time() - cached.get("updated_at", 0) > self.pipeline_cache_ttl:
            self.pipeline_cache.pop(doc_id, None)
            return None

        return cached

    def _touch_pipeline_cache(self, doc_id: Optional[str], **kwargs) -> Optional[Dict[str, Any]]:
        """写入文档处理缓存"""
        if not doc_id:
            return None

        cached = self.pipeline_cache.get(doc_id, {})
        cached.update(kwargs)
        cached["updated_at"] = time.time()
        self.pipeline_cache[doc_id] = cached
        return cached

    async def process_document_pipeline_async(
        self,
        file_path: str,
        doc_id: Optional[str] = None,
        include_triplets: bool = False,
        force_refresh: bool = False
    ) -> Dict[str, Any]:
        """异步文档处理流水线，返回 text/chunks/triplets"""
        cached = None if force_refresh else self._get_pipeline_cache(doc_id)
        has_required_stage = cached and (
            cached.get("triplets") is not None if include_triplets else cached.get("chunks") is not None
        )

        if has_required_stage:
            return {
                "text": cached.get("text", ""),
                "chunks": cached.get("chunks", []),
                "triplets": cached.get("triplets", [])
            }

        text = cached.get("text") if cached and cached.get("text") is not None else self.parse_document(file_path)
        chunks = cached.get("chunks") if cached and cached.get("chunks") is not None else self.split_text(text)
        triplets = cached.get("triplets") if cached and cached.get("triplets") is not None else None

        if include_triplets and triplets is None:
            triplets = await self.extract_batch_async(chunks)

        self._touch_pipeline_cache(doc_id, text=text, chunks=chunks, triplets=triplets)
        return {
            "text": text,
            "chunks": chunks,
            "triplets": triplets or []
        }

    def get_qa_client(self, model_name: Optional[str] = None):
        """获取AI问答客户端和模型名称"""
        if model_name is None:
            model_name = DEFAULT_QA_MODEL

        if model_name not in self.qa_clients:
            raise ValueError(f"模型 '{model_name}' 不存在。可用模型: {list(self.qa_clients.keys())}")

        return self.qa_clients[model_name]["client"], self.qa_clients[model_name]["model"]

    # ==================== 文档处理 ====================

    def parse_pdf(self, file_path: str) -> str:
        """解析PDF文档"""
        try:
            reader = PdfReader(file_path)
            text = ""
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
            return text
        except Exception as e:
            raise Exception(f"PDF解析失败: {str(e)}")

    def parse_document(self, file_path: str) -> str:
        """解析文档（支持PDF/TXT/DOCX/PPTX）"""
        ext = os.path.splitext(file_path)[1].lower()

        if ext == '.pdf':
            return self.parse_pdf(file_path)
        elif ext == '.txt':
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        elif ext == '.docx':
            from docx import Document
            doc = Document(file_path)
            text_parts = []
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            text_parts.append(cell.text)
            return '\n'.join(text_parts)
        elif ext == '.pptx':
            from pptx import Presentation
            prs = Presentation(file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return '\n'.join(text)
        else:
            raise Exception(f"不支持的文件格式: {ext}")

    def split_text(self, text: str) -> List[Dict[str, Any]]:
        """分块文本"""
        chunks = self.text_splitter.split_text(text)
        chunk_list = []
        for idx, chunk in enumerate(chunks):
            chunk_list.append({
                "index": idx,
                "content": chunk,
                "length": len(chunk)
            })
        return chunk_list

    # ==================== 实体关系抽取 ====================

    def extract_entities_and_relations(self, text: str) -> List[Dict[str, str]]:
        """使用LLM提取实体和关系（三元组）"""
        prompt = f"""你是一个知识图谱专家。请从以下文本中提取实体和关系，以JSON格式输出。

要求：
1. 识别文本中的关键实体（人物、组织、地点、概念等）
2. 识别实体之间的关系
3. 输出格式为JSON数组，每个元素包含：head（头实体）、relation（关系）、tail（尾实体）、head_type（头实体类型）、tail_type（尾实体类型）

文本内容：
{text}

请直接返回JSON数组，不要有其他文字说明："""

        try:
            response = self.extraction_client.chat.completions.create(
                model=EXTRACTION_MODEL,
                messages=[
                    {"role": "system", "content": "你是一个专业的知识图谱构建助手。"},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.1,
                max_tokens=8000
            )

            if response.choices and len(response.choices) > 0:
                content = response.choices[0].message.content
                
                if content is None:
                    reasoning_content = getattr(response.choices[0].message, 'reasoning_content', None)
                    if reasoning_content:
                        content = reasoning_content
                        logging.warning(f"[EXTRACTION] content 为 None，从 reasoning_content 获取内容")
                    else:
                        logging.error(f"[EXTRACTION] content 和 reasoning_content 都为 None")
                        return []
                
                result_text = content.strip()
            else:
                logging.error(f"[EXTRACTION] 没有有效的 choices")
                return []

            if "```json" in result_text:
                result_text = result_text.split("```json")[1].split("```")[0].strip()
            elif "```" in result_text:
                result_text = result_text.split("```")[1].split("```")[0].strip()

            triplets = json.loads(result_text)
            return triplets if isinstance(triplets, list) else []

        except Exception as e:
            logging.error(f"[EXTRACTION] 实体关系抽取失败: {str(e)}")
            return []

    async def extract_batch_async(self, chunks: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        """并发处理多个文本块"""
        semaphore = asyncio.Semaphore(8)

        async def sem_task(chunk):
            async with semaphore:
                loop = asyncio.get_event_loop()
                await asyncio.sleep(0.5)
                try:
                    result = await loop.run_in_executor(
                        self.executor,
                        self.extract_entities_and_relations,
                        chunk["content"]
                    )
                    return result
                except Exception as e:
                    print(f"抽取块 {chunk['index']} 失败: {e}")
                    return []

        tasks = [sem_task(chunk) for chunk in chunks]
        results = await asyncio.gather(*tasks)

        all_triplets = []
        for idx, triplets_list in enumerate(results):
            if triplets_list and isinstance(triplets_list, list):
                for triplet in triplets_list:
                    triplet["source_chunk_index"] = chunks[idx]["index"]
                    all_triplets.append(triplet)

        return all_triplets

    # ==================== Neo4j图谱存储 ====================

    def save_triplets_to_neo4j(self, triplets: List[Dict[str, Any]], doc_id: str):
        """将三元组保存到Neo4j"""
        with self.neo4j_driver.session() as session:
            session.run("MERGE (d:Document {id: $doc_id})", doc_id=doc_id)

            for triplet in triplets:
                try:
                    head = triplet.get("head", "")
                    tail = triplet.get("tail", "")
                    relation = triplet.get("relation", "RELATED_TO")
                    head_type = triplet.get("head_type", "Entity")
                    tail_type = triplet.get("tail_type", "Entity")

                    if not head or not tail:
                        continue

                    safe_head_type = "".join([c if c.isalnum() or c == "_" else "_" for c in head_type]) if head_type else "Entity"
                    safe_tail_type = "".join([c if c.isalnum() or c == "_" else "_" for c in tail_type]) if tail_type else "Entity"

                    if not safe_head_type or safe_head_type[0].isdigit():
                        safe_head_type = "Entity"
                    if not safe_tail_type or safe_tail_type[0].isdigit():
                        safe_tail_type = "Entity"

                    session.run(
                        f"MERGE (h:{safe_head_type} {{name: $head}}) "
                        "ON CREATE SET h.created_at = timestamp() "
                        "MERGE (d:Document {id: $doc_id}) "
                        "MERGE (h)-[:FROM_DOCUMENT]->(d)",
                        head=head, doc_id=doc_id
                    )

                    session.run(
                        f"MERGE (t:{safe_tail_type} {{name: $tail}}) "
                        "ON CREATE SET t.created_at = timestamp() "
                        "MERGE (d:Document {id: $doc_id}) "
                        "MERGE (t)-[:FROM_DOCUMENT]->(d)",
                        tail=tail, doc_id=doc_id
                    )

                    safe_relation = "".join([c if c.isalnum() or c == "_" else "_" for c in relation])
                    session.run(
                        f"MATCH (h {{name: $head}}) "
                        f"MATCH (t {{name: $tail}}) "
                        f"MERGE (h)-[r:{safe_relation}]->(t) "
                        "ON CREATE SET r.created_at = timestamp()",
                        head=head, tail=tail
                    )

                except Exception as e:
                    print(f"保存三元组失败: {triplet}, 错误: {str(e)}")
                    continue

    def get_graph_data(self, doc_id: Optional[str] = None, limit: int = 100) -> Dict[str, Any]:
        """获取图谱数据"""
        with self.neo4j_driver.session() as session:
            if doc_id:
                result = session.run(
                    """
                    MATCH (n)-[r]->(m)
                    WHERE type(r) <> 'FROM_DOCUMENT'
                      AND NOT 'Document' IN labels(n)
                      AND NOT 'Document' IN labels(m)
                      AND (
                        (n)-[:FROM_DOCUMENT]->(:Document {id: $doc_id})
                        OR (m)-[:FROM_DOCUMENT]->(:Document {id: $doc_id})
                      )
                    RETURN n, r, m
                    LIMIT $limit
                    """,
                    doc_id=doc_id, limit=limit
                )
            else:
                result = session.run(
                    """
                    MATCH (n)-[r]->(m)
                    WHERE NOT type(r) = 'FROM_DOCUMENT'
                    RETURN n, r, m
                    LIMIT $limit
                    """,
                    limit=limit
                )

            nodes = {}
            edges = []

            for record in result:
                n = record["n"]
                r = record["r"]
                m = record["m"]

                if n.element_id not in nodes:
                    nodes[n.element_id] = {
                        "id": n.element_id,
                        "label": n.get("name", "Unknown"),
                        "type": list(n.labels)[0] if n.labels else "Entity"
                    }

                if m.element_id not in nodes:
                    nodes[m.element_id] = {
                        "id": m.element_id,
                        "label": m.get("name", "Unknown"),
                        "type": list(m.labels)[0] if m.labels else "Entity"
                    }

                edges.append({
                    "source": n.element_id,
                    "target": m.element_id,
                    "label": r.type
                })

            return {"nodes": list(nodes.values()), "edges": edges}

    # ==================== ChromaDB向量存储 ====================

    def save_chunks_to_chromadb(self, chunks: List[Dict[str, Any]], doc_id: str):
        """将文本块保存到ChromaDB"""
        texts = [chunk["content"] for chunk in chunks]
        logging.info(f"[VECTOR] 开始处理 {len(texts)} 个文本块")

        try:
            logging.info(f"[VECTOR] 正在生成 embeddings...")
            embeddings_response = self.embed_client.embeddings.create(
                model=EMBED_MODEL,
                input=texts,
                encoding_format="float"
            )
            embeddings = [item.embedding for item in embeddings_response.data]
            logging.info(f"[VECTOR] ✅ 成功生成 {len(embeddings)} 个 embeddings")

            BATCH_SIZE = 50
            total_chunks = len(chunks)
            total_batches = (total_chunks + BATCH_SIZE - 1) // BATCH_SIZE

            saved_count = 0
            
            for batch_idx in range(total_batches):
                start_idx = batch_idx * BATCH_SIZE
                end_idx = min(start_idx + BATCH_SIZE, total_chunks)
                
                batch_chunks = chunks[start_idx:end_idx]
                batch_texts = texts[start_idx:end_idx]
                batch_embeddings = embeddings[start_idx:end_idx]
                
                batch_ids = [f"{doc_id}_{chunk['index']}" for chunk in batch_chunks]
                batch_metadatas = [
                    {"doc_id": doc_id, "chunk_index": chunk["index"], "length": chunk["length"]}
                    for chunk in batch_chunks
                ]

                try:
                    self.collection.add(
                        ids=batch_ids,
                        embeddings=batch_embeddings,
                        documents=batch_texts,
                        metadatas=batch_metadatas
                    )
                    saved_count += len(batch_chunks)
                    logging.info(f"[VECTOR] ✅ 第 {batch_idx + 1}/{total_batches} 批保存成功")
                except Exception as batch_error:
                    logging.error(f"[VECTOR] ❌ 第 {batch_idx + 1}/{total_batches} 批保存失败: {batch_error}")
                    continue

            if saved_count == 0:
                raise Exception("所有批次保存失败")

        except Exception as e:
            logging.error(f"[VECTOR] ❌ 向量存储失败: {str(e)}")
            raise Exception(f"向量存储失败: {str(e)}")

    # ==================== Milvus 向量存储 ====================

    def save_chunks_to_milvus(self, chunks: List[Dict[str, Any]], doc_id: str):
        """将文本块保存到 Milvus（高性能向量存储）"""
        if not self.milvus_client or not self.milvus_client.connected:
            logging.warning("[MILVUS] Milvus 未连接，回退到 ChromaDB")
            return self.save_chunks_to_chromadb(chunks, doc_id)
        
        texts = [chunk["content"] for chunk in chunks]
        logging.info(f"[MILVUS] 开始处理 {len(texts)} 个文本块")

        try:
            logging.info(f"[MILVUS] 正在生成 embeddings...")
            embeddings_response = self.embed_client.embeddings.create(
                model=EMBED_MODEL,
                input=texts,
                encoding_format="float"
            )
            embeddings = [item.embedding for item in embeddings_response.data]
            logging.info(f"[MILVUS] ✅ 成功生成 {len(embeddings)} 个 embeddings")

            # 准备数据
            ids = [f"{doc_id}_{chunk['index']}" for chunk in chunks]
            metadatas = [
                {"doc_id": doc_id, "chunk_index": chunk["index"], "length": chunk["length"]}
                for chunk in chunks
            ]

            # 批量插入
            BATCH_SIZE = 100
            total_chunks = len(chunks)
            total_batches = (total_chunks + BATCH_SIZE - 1) // BATCH_SIZE

            saved_count = 0
            for batch_idx in range(total_batches):
                start_idx = batch_idx * BATCH_SIZE
                end_idx = min(start_idx + BATCH_SIZE, total_chunks)

                batch_ids = ids[start_idx:end_idx]
                batch_embeddings = embeddings[start_idx:end_idx]
                batch_texts = texts[start_idx:end_idx]
                batch_metadatas = metadatas[start_idx:end_idx]

                try:
                    success = self.milvus_client.insert(
                        ids=batch_ids,
                        vectors=batch_embeddings,
                        documents=batch_texts,
                        metadatas=batch_metadatas
                    )
                    if success:
                        saved_count += len(batch_ids)
                        logging.info(f"[MILVUS] ✅ 第 {batch_idx + 1}/{total_batches} 批保存成功")
                except Exception as batch_error:
                    logging.error(f"[MILVUS] ❌ 第 {batch_idx + 1}/{total_batches} 批保存失败: {batch_error}")
                    continue

            if saved_count == 0:
                raise Exception("所有批次保存失败")
            logging.info(f"[MILVUS] ✅ 成功保存 {saved_count} 条数据")

        except Exception as e:
            logging.error(f"[MILVUS] ❌ 向量存储失败: {str(e)}")
            raise Exception(f"Milvus 向量存储失败: {str(e)}")

    def _search_with_milvus(self, query_embedding: List[float], top_k: int) -> List[Dict[str, Any]]:
        """使用 Milvus 进行向量检索"""
        if not self.milvus_client or not self.milvus_client.connected:
            logging.warning("[MILVUS] Milvus 未连接，回退到 ChromaDB")
            return []

        try:
            results = self.milvus_client.search(
                query_vectors=[query_embedding],
                top_k=top_k
            )

            chunks = []
            if results and len(results) > 0:
                for hit in results[0]:
                    # Milvus 返回的是距离，需要转换为相似度
                    # 距离越小越相似
                    distance = hit.get("distance", 1.0)
                    # 对于 COSINE/IP 距离，越大越相似
                    # 对于 L2 距离，越小越相似
                    if distance >= 0:
                        # 假设使用 COSINE，distance 就是相似度
                        similarity = distance
                        # 转换 distance（使其与其他代码兼容）
                        vector_distance = 1 - similarity
                    else:
                        vector_distance = abs(distance)
                        similarity = 1 / (1 + vector_distance)

                    # 过滤相似度
                    if vector_distance < VECTOR_SIMILARITY_THRESHOLD:
                        chunks.append({
                            "content": hit.get("content", ""),
                            "metadata": {
                                "doc_id": hit.get("doc_id", ""),
                                "chunk_index": hit.get("chunk_index", 0),
                                "length": hit.get("length", 0)
                            },
                            "distance": vector_distance,
                            "similarity": similarity
                        })

            return chunks

        except Exception as e:
            logging.error(f"[MILVUS] 检索失败: {e}")
            return []

    # ==================== Query Rewrite ====================

    def rewrite_query(self, query: str) -> List[str]:
        """
        使用 LLM 改写用户查询，生成多个等价查询以提升检索召回率

        Args:
            query: 用户原始查询

        Returns:
            改写后的查询列表
        """
        if not QUERY_REWRITE_ENABLED:
            return [query]

        prompt = f"""你是一个查询改写专家。请将用户的问题改写成 {QUERY_REWRITE_COUNT} 个不同的搜索查询。

要求：
1. 生成正好 {QUERY_REWRITE_COUNT} 个改写后的查询
2. 每个查询应该能检索到不同的相关文档
3. 包含同义词改写、角度转换、简化表达等
4. 输出格式为 JSON 数组，不要有其他文字说明
5. 保留中文和英文的混合表达

用户问题：{query}

请直接返回 JSON 数组，例如：["docker 部署步骤", "Docker 容器部署教程", "如何使用 docker compose 部署"]"""

        try:
            response = self.extraction_client.chat.completions.create(
                model=EXTRACTION_MODEL,
                messages=[
                    {"role": "system", "content": "你是一个专业的查询改写助手。"},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.3,
                max_tokens=1000
            )

            if response.choices and len(response.choices) > 0:
                content = response.choices[0].message.content
                if content is None:
                    reasoning_content = getattr(response.choices[0].message, 'reasoning_content', None)
                    if reasoning_content:
                        content = reasoning_content
                    else:
                        return [query]

                result_text = content.strip()

                # 解析 JSON 数组
                if "```json" in result_text:
                    result_text = result_text.split("```json")[1].split("```")[0].strip()
                elif "```" in result_text:
                    result_text = result_text.split("```")[1].split("```")[0].strip()

                rewritten_queries = json.loads(result_text)
                if isinstance(rewritten_queries, list) and len(rewritten_queries) >= 1:
                    print(f"[QUERY_REWRITE] ✅ 成功改写为 {len(rewritten_queries)} 个查询")
                    return rewritten_queries[:QUERY_REWRITE_COUNT]
                else:
                    return [query]
            else:
                return [query]

        except Exception as e:
            print(f"[QUERY_REWRITE] ❌ 改写失败: {e}")
            return [query]

    def _search_single_query(self, query: str, n_results: int) -> List[Dict[str, Any]]:
        """
        单个查询的向量检索（内部方法，不走缓存）

        Args:
            query: 查询文本
            n_results: 返回结果数量

        Returns:
            检索结果列表
        """
        try:
            # 生成查询向量
            query_embedding_response = self.embed_client.embeddings.create(
                model=EMBED_MODEL,
                input=[query],
                encoding_format="float"
            )
            query_embedding = query_embedding_response.data[0].embedding

            # 确定检索数量：如果启用 Reranker，先检索更多结果
            search_n = RERANKER_INITIAL_TOP_K if (self.reranker and RERANKER_ENABLED) else n_results

            # 检索
            results = self.collection.query(
                query_embeddings=[query_embedding],
                n_results=search_n
            )

            chunks = []
            # 格式化结果并过滤相似度
            if results and "documents" in results and results["documents"]:
                for i in range(len(results["documents"][0])):
                    try:
                        distance = results["distances"][0][i] if "distances" in results else None
                        if distance is not None and distance < VECTOR_SIMILARITY_THRESHOLD:
                            chunks.append({
                                "content": results["documents"][0][i],
                                "metadata": results["metadatas"][0][i] if "metadatas" in results else {},
                                "distance": distance,
                                "query": query  # 记录来自哪个查询
                            })
                    except (IndexError, KeyError):
                        continue

                # Reranker 重排序
                if self.reranker and RERANKER_ENABLED and len(chunks) > 1:
                    try:
                        candidate_contents = [
                            c["content"][:RERANKER_MAX_LENGTH] if len(c["content"]) > RERANKER_MAX_LENGTH
                            else c["content"] for c in chunks
                        ]
                        rerank_scores = self.reranker.score(query, candidate_contents)

                        scored_chunks = []
                        for idx, c in enumerate(chunks):
                            rerank_score = rerank_scores[idx] if idx < len(rerank_scores) else 0
                            vector_distance = c.get("distance", 1.0)
                            final_score = (RERANKER_WEIGHT * float(rerank_score) + VECTOR_WEIGHT * (1 - vector_distance))

                            if float(rerank_score) >= RERANKER_SCORE_THRESHOLD:
                                scored_chunks.append({
                                    **c,
                                    "rerank_score": float(rerank_score),
                                    "final_score": final_score
                                })

                        scored_chunks.sort(key=lambda x: x.get("final_score", 0), reverse=True)
                        chunks = scored_chunks[:n_results]
                    except Exception:
                        chunks = chunks[:n_results]
                else:
                    chunks = chunks[:n_results]

            return chunks

        except Exception as e:
            print(f"[QUERY_REWRITE] 单查询检索失败: {e}")
            return []

    def _merge_search_results(self, all_results: List[Dict[str, Any]], top_k: int) -> List[Dict[str, Any]]:
        """
        合并多个查询的检索结果并进行去重

        Args:
            all_results: 所有查询的检索结果
            top_k: 最终返回的结果数量

        Returns:
            合并去重后的结果列表
        """
        if not all_results:
            return []

        # 按 final_score（如果有）或 distance 排序
        def get_score(item):
            if "final_score" in item:
                return item["final_score"]
            elif "rerank_score" in item:
                return item["rerank_score"]
            else:
                return -item.get("distance", 0)

        all_results.sort(key=get_score, reverse=True)

        # 去重：基于 content 相似度去重
        merged = []
        seen_contents = set()

        for item in all_results:
            content = item.get("content", "")
            # 简单的去重：完全相同的内容去重
            if content not in seen_contents:
                seen_contents.add(content)
                merged.append(item)

            if len(merged) >= top_k:
                break

        return merged

    # ==================== 向量检索（带缓存和 Query Rewrite）====================

    def search_similar_chunks(self, query: str, n_results: int = None) -> List[Dict[str, Any]]:
        """向量检索相似文本块（支持 Query Rewrite 多查询检索、缓存、Reranking重排序）"""

        if n_results is None:
            n_results = VECTOR_SEARCH_TOP_K

        logger = logging.getLogger(__name__)

        # 🔥 检查缓存（优先返回）
        cached_results = get_search_cache(query)
        if cached_results is not None:
            logger.info(f'[CACHE] 直接返回缓存结果: {len(cached_results)} 条')
            return cached_results

        chunks = []
        query_embedding = None

        try:
            start_ts = time.time()
            logger.info('[VECTOR] search_similar_chunks start')

            # 生成查询向量（统一生成，用于检索和缓存）
            query_embedding_response = self.embed_client.embeddings.create(
                model=EMBED_MODEL,
                input=[query],
                encoding_format="float"
            )
            query_embedding = query_embedding_response.data[0].embedding

            # 判断使用 Milvus 还是 ChromaDB
            use_milvus = self.milvus_client and self.milvus_client.connected
            
            # 1️⃣ Query Rewrite：生成多个等价查询
            if QUERY_REWRITE_ENABLED:
                logger.info('[QUERY_REWRITE] 开始改写查询...')
                rewritten_queries = self.rewrite_query(query)
                logger.info(f'[QUERY_REWRITE] 改写为 {len(rewritten_queries)} 个查询: {rewritten_queries[:3]}...')

                # 2️⃣ 并行检索所有改写后的查询
                all_results = []
                search_top_k = max(QUERY_REWRITE_TOP_K, n_results)  # 每个查询检索足够多的结果

                for q in rewritten_queries:
                    if use_milvus:
                        # 使用 Milvus 检索
                        q_emb_response = self.embed_client.embeddings.create(
                            model=EMBED_MODEL,
                            input=[q],
                            encoding_format="float"
                        )
                        q_embedding = q_emb_response.data[0].embedding
                        results = self._search_with_milvus(q_embedding, search_top_k)
                    else:
                        # 使用 ChromaDB 检索
                        results = self._search_single_query(q, search_top_k)
                    all_results.extend(results)
                    logger.info(f'[QUERY_REWRITE] 查询 "{q[:20]}..." 检索到 {len(results)} 条结果')

                # 3️⃣ 合并去重结果
                chunks = self._merge_search_results(all_results, QUERY_REWRITE_MERGE_TOP_K)
                logger.info(f'[QUERY_REWRITE] 合并后共 {len(chunks)} 条结果')
            else:
                # 直接检索原始查询
                if use_milvus:
                    # 使用 Milvus 检索
                    chunks = self._search_with_milvus(query_embedding, n_results)
                else:
                    # 使用 ChromaDB 检索
                    search_n = RERANKER_INITIAL_TOP_K if (self.reranker and RERANKER_ENABLED) else n_results
                    results = self.collection.query(
                        query_embeddings=[query_embedding],
                        n_results=search_n
                    )

                    # 格式化结果并过滤相似度
                    if results and "documents" in results and results["documents"]:
                        candidates = []
                        for i in range(len(results["documents"][0])):
                            try:
                                distance = results["distances"][0][i] if "distances" in results else None
                                if distance is not None and distance < VECTOR_SIMILARITY_THRESHOLD:
                                    candidates.append({
                                        "content": results["documents"][0][i],
                                        "metadata": results["metadatas"][0][i] if "metadatas" in results else {},
                                        "distance": distance
                                    })
                            except (IndexError, KeyError) as e:
                                logger.warning(f'[VECTOR] 处理结果索引 {i} 时出错: {e}')
                                continue

                        # Reranker 重排序
                        if self.reranker and RERANKER_ENABLED and len(candidates) > 1:
                            try:
                                candidate_contents = [
                                    c["content"][:RERANKER_MAX_LENGTH] if len(c["content"]) > RERANKER_MAX_LENGTH
                                    else c["content"] for c in candidates
                                ]
                                rerank_scores = self.reranker.score(query, candidate_contents)

                                scored_candidates = []
                                for idx, c in enumerate(candidates):
                                    rerank_score = rerank_scores[idx] if idx < len(rerank_scores) else 0
                                    vector_distance = c.get("distance", 1.0)
                                    final_score = (RERANKER_WEIGHT * float(rerank_score) + VECTOR_WEIGHT * (1 - vector_distance))

                                    if float(rerank_score) >= RERANKER_SCORE_THRESHOLD:
                                        scored_candidates.append({
                                            **c,
                                            "rerank_score": float(rerank_score),
                                            "final_score": final_score
                                        })

                                scored_candidates.sort(key=lambda x: x.get("final_score", 0), reverse=True)
                                chunks = scored_candidates[:n_results]
                            except Exception as rerank_err:
                                logger.warning(f'[RERANKER] 重排序失败: {rerank_err}')
                                chunks = candidates[:n_results]
                        else:
                            chunks = candidates[:n_results]

            elapsed_ms = int((time.time() - start_ts) * 1000)
            logger.info(f'[VECTOR] search_similar_chunks done elapsed_ms={elapsed_ms} count={len(chunks)}')

            # 💾 保存结果到缓存（传入 embedding 启用语义缓存）
            if chunks and query_embedding is not None:
                set_search_cache(query, chunks, query_embedding)

            return chunks

        except Exception as e:
            logger.error(f'[VECTOR] search_similar_chunks exception: {str(e)}')
            print(f"向量检索失败: {str(e)}")
            return []

    # ==================== 图检索 ====================

    def search_relevant_chunks(self, query: str, n_results: int = None) -> List[Dict[str, Any]]:
        """统一 RAG 检索入口：优先使用混合检索，否则回退到单一检索方法"""
        if n_results is None:
            n_results = VECTOR_SEARCH_TOP_K

        if DISABLE_VECTOR_SEARCH:
            if BM25_ENABLED and ES_AVAILABLE:
                return search_bm25_chunks(query, top_k=n_results)
            return []

        if BM25_ENABLED and ES_AVAILABLE:
            vector_results = self.search_similar_chunks(query, n_results)
            bm25_results = search_bm25_chunks(query, top_k=n_results)

            if vector_results and bm25_results:
                return rrf_fusion([vector_results, bm25_results], k=RRF_K)[:n_results]
            if vector_results:
                return vector_results[:n_results]
            return bm25_results[:n_results]

        return self.search_similar_chunks(query, n_results)

    def search_graph_neighbors(self, entities: List[str], hops: int = None) -> Dict[str, Any]:
        """图检索：查找实体的邻居节点"""
        if not entities:
            return {"nodes": [], "edges": []}

        if hops is None:
            hops = GRAPH_SEARCH_HOPS

        with self.neo4j_driver.session() as session:
            query = f"""
            MATCH p = (n)-[r*1..{hops}]-(m)
            WHERE n.name IN $entities
            AND NOT 'Document' IN labels(n)
            AND NOT 'Document' IN labels(m)
            AND ALL(rel IN r WHERE type(rel) <> 'FROM_DOCUMENT')
            RETURN n, r as r_list, m
            LIMIT 50
            """

            result = session.run(query, entities=entities)

            nodes = {}
            edges = []

            for record in result:
                n = record["n"]
                m = record["m"]
                r_list = record["r_list"]

                if n.element_id not in nodes:
                    nodes[n.element_id] = {
                        "id": str(n.element_id),
                        "label": str(n.get("name", "Unknown")),
                        "type": str(list(n.labels)[0]) if n.labels else "Entity",
                        "highlighted": True
                    }

                if m.element_id not in nodes:
                    nodes[m.element_id] = {
                        "id": str(m.element_id),
                        "label": str(m.get("name", "Unknown")),
                        "type": str(list(m.labels)[0]) if m.labels else "Entity",
                        "highlighted": True
                    }

                for rel in r_list:
                    edge_id = str(rel.element_id)
                    if not any(e['id'] == edge_id for e in edges):
                        edges.append({
                            "id": edge_id,
                            "source": str(rel.start_node.element_id),
                            "target": str(rel.end_node.element_id),
                            "label": str(rel.type)
                        })

            return {"nodes": list(nodes.values()), "edges": edges}

    def extract_entities_from_question(self, question: str) -> List[str]:
        """从问题中提取实体"""
        logging.info(f'[EXTRACTION] Extract entities from question: {question}')

        prompt = f"""从以下问题中提取关键实体名称，以JSON数组格式返回。

问题：{question}

只返回实体名称数组，例如：["张三", "阿里巴巴"]

请直接返回JSON数组："""

        try:
            response = self.extraction_client.chat.completions.create(
                model=EXTRACTION_MODEL,
                messages=[
                    {"role": "system", "content": "你是一个实体识别助手。"},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.1,
                max_tokens=500,
                timeout=10.0
            )

            if response.choices and len(response.choices) > 0:
                content = response.choices[0].message.content
                
                if content is None:
                    reasoning_content = getattr(response.choices[0].message, 'reasoning_content', None)
                    if reasoning_content:
                        content = reasoning_content
                    else:
                        return []
                
                result_text = content.strip()
            else:
                return []

            if "```json" in result_text:
                result_text = result_text.split("```json")[1].split("```")[0].strip()
            elif "```" in result_text:
                result_text = result_text.split("```")[1].split("```")[0].strip()

            entities = json.loads(result_text)
            return entities if isinstance(entities, list) else []

        except Exception as e:
            logging.error(f"[EXTRACTION] 实体提取失败: {str(e)}")
            return []

    # ==================== RAG问答 ====================

    async def answer_question(self, question: str, stream: bool = False, model_name: Optional[str] = None) -> Dict[str, Any]:
        """RAG问答（混合检索 + LLM生成）"""
        qa_client, qa_model = self.get_qa_client(model_name)

        vector_chunks = self.search_relevant_chunks(question, n_results=5)

        entities = self.extract_entities_from_question(question)
        graph_data = self.search_graph_neighbors(entities, hops=2)

        context_parts = []

        if vector_chunks:
            context_parts.append("【相关文档片段】")
            for i, chunk in enumerate(vector_chunks, 1):
                context_parts.append(f"{i}. {chunk['content']}")

        if graph_data["nodes"]:
            context_parts.append("\n【相关知识图谱】")
            for node in graph_data["nodes"][:10]:
                context_parts.append(f"- {node['label']} ({node['type']})")

        if not context_parts:
            prompt = f"""请回答以下问题。注意：未找到相关文档或知识图谱信息，请基于你的知识直接回答。

问题：{question}

请提供详细的答案："""
        else:
            context = "\n".join(context_parts)
            prompt = f"""请基于以下上下文回答问题。如果上下文中没有相关信息，请如实说明。

上下文：
{context}

问题：{question}

请提供详细的答案："""

        try:
            if stream:
                response = qa_client.chat.completions.create(
                    model=qa_model,
                    messages=[
                        {"role": "system", "content": "你是一个专业的知识问答助手。"},
                        {"role": "user", "content": prompt}
                    ],
                    temperature=0.3,
                    max_tokens=2000,
                    stream=True
                )
                return {"stream": response, "sources": {"vector_chunks": vector_chunks, "graph_data": graph_data}}
            else:
                response = qa_client.chat.completions.create(
                    model=qa_model,
                    messages=[
                        {"role": "system", "content": "你是一个专业的知识问答助手。"},
                        {"role": "user", "content": prompt}
                    ],
                    temperature=0.3,
                    max_tokens=2000
                )
                answer = response.choices[0].message.content
                return {"answer": answer, "sources": {"vector_chunks": vector_chunks, "graph_data": graph_data}}

        except Exception as e:
            raise Exception(f"答案生成失败: {str(e)}")

    # ==================== 其他方法（简化展示）====================

    def call_qidian_api_stream(self, prompt: str, session_id: str, satoken: str, model_id: int = 32):
        """调用 qidian API 流式问答"""
        url = "https://qidianai.xyz/api/user/chat/send"
        headers = {
            "Accept": "text/event-stream",
            "Content-Type": "application/json",
            "satoken": satoken,
            "Cookie": f"satoken={satoken}"
        }
        payload = {"modelId": model_id, "content": prompt, "sessionId": session_id, "attachments": []}

        try:
            with requests.post(url, headers=headers, json=payload, stream=True, timeout=300) as resp:
                if resp.status_code >= 400:
                    raise Exception(f"Qidian API 返回 {resp.status_code} 错误")
                for block in self._parse_sse_blocks(resp):
                    event_type, data = self._parse_sse_block(block)
                    text = self._decode_data_text(data)
                    if text and text.strip() not in ("[DONE]", "DONE"):
                        if event_type in (None, "message", "start"):
                            if event_type != "start" and text:
                                yield text
        except Exception as e:
            raise Exception(f"Qidian API 请求失败: {str(e)}")

    def _parse_sse_blocks(self, resp):
        block = []
        for raw in resp.iter_lines(decode_unicode=True):
            if raw is None:
                continue
            line = raw.rstrip("\r")
            if line == "":
                if block:
                    yield block
                    block = []
                continue
            block.append(line)
        if block:
            yield block

    def _parse_sse_block(self, lines):
        event_type = None
        data_lines = []
        for line in lines:
            if line.startswith("event:"):
                event_type = line.split(":", 1)[1].strip()
            elif line.startswith("data:"):
                data_lines.append(line.split(":", 1)[1].lstrip())
        return event_type, "\n".join(data_lines)

    def _decode_data_text(self, data):
        s = (data or "").strip()
        if not s or s in ("[DONE]", "DONE"):
            return s
        if (s.startswith("{") and s.endswith("}")) or (s.startswith("[") and s.endswith("]")):
            try:
                obj = json.loads(s)
                return self._extract_text_from_json(obj) or s
            except Exception:
                return data
        return data

    def _extract_text_from_json(self, obj):
        if obj is None:
            return None
        if isinstance(obj, str):
            return obj
        if isinstance(obj, (int, float, bool)):
            return str(obj)
        if isinstance(obj, dict):
            for key in ("content", "message", "text", "delta", "answer", "data"):
                if key in obj:
                    val = obj[key]
                    if isinstance(val, (dict, list)):
                        extracted = self._extract_text_from_json(val)
                        if extracted:
                            return extracted
                    elif val is not None:
                        return str(val)
            if "choices" in obj and isinstance(obj["choices"], list) and obj["choices"]:
                return self._extract_text_from_json(obj["choices"][0])
        if isinstance(obj, list):
            parts = []
            for item in obj:
                t = self._extract_text_from_json(item)
                if t:
                    parts.append(t)
            return "".join(parts) if parts else None
        return None

    def call_yansd_api_stream(self, prompt: str, model: str, session_cookie: str, new_api_user: str):
        """调用 yansd API 流式问答"""
        url = "https://yansd666.com/pg/chat/completions"
        headers = {
            "Content-Type": "application/json",
            "Accept": "text/event-stream",
            "Cookie": session_cookie,
            "new-api-user": new_api_user
        }
        payload = {
            "model": model,
            "group": "default",
            "messages": [{"role": "user", "content": prompt}],
            "stream": True
        }

        try:
            with requests.post(url, headers=headers, json=payload, stream=True, timeout=300) as resp:
                if resp.status_code >= 400:
                    raise Exception(f"Yansd API 返回 {resp.status_code} 错误")
                resp.encoding = 'utf-8'
                for line in resp.iter_lines(decode_unicode=True):
                    if not line or not line.startswith("data: "):
                        continue
                    data_str = line[6:].strip()
                    if data_str == "[DONE]":
                        break
                    try:
                        chunk = json.loads(data_str)
                        content = chunk["choices"][0]["delta"].get("content")
                        if content:
                            yield content
                    except (json.JSONDecodeError, KeyError, IndexError):
                        continue
        except Exception as e:
            raise Exception(f"Yansd API 请求失败: {str(e)}")

    # ==================== 并行流式RAG问答 ====================

    async def answer_question_parallel_stream(self, question: str, conversation_id: Optional[str] = None, model_name: Optional[str] = None):
        """并行流式RAG问答"""
        import logging
        from redis_utils import get_conversation_history, save_conversation_message
        
        logger = logging.getLogger(__name__)
        logger.info(f"🟢 answer_question_parallel_stream 开始: {question[:50]}...")

        try:
            qa_client, qa_model = self.get_qa_client(model_name)
        except Exception as e:
            logger.error(f"❌ 获取QA客户端失败: {e}")
            raise

        vector_chunks = []
        graph_data = {"nodes": [], "edges": []}
        full_answer = ""

        # 获取对话历史
        history = []
        if conversation_id:
            try:
                history = get_conversation_history(conversation_id)
            except Exception as e:
                logger.error(f"获取对话历史失败: {e}")

        # 向量检索任务
        async def vector_search_task():
            logger.info(f"🔍 [向量检索] 开始")
            try:
                chunks = self.search_relevant_chunks(question, 5)
                vector_chunks.extend(chunks)
                logger.info(f"✅ [向量检索] 完成，找到 {len(chunks)} 个相似文档")
                return {"type": "vector_chunks", "data": chunks}
            except Exception as e:
                logger.error(f"❌ [向量检索] 失败: {e}")
                return {"type": "vector_chunks", "data": [], "error": str(e)}

        # 图检索任务
        async def graph_search_task():
            logger.info(f"🔍 [图检索] 开始")
            try:
                loop = asyncio.get_event_loop()
                entities = await loop.run_in_executor(self.executor, self.extract_entities_from_question, question)
                graph = await loop.run_in_executor(self.executor, self.search_graph_neighbors, entities, 2)
                graph_data.update(graph)
                logger.info(f"✅ [图检索] 完成，找到 {len(graph.get('nodes', []))} 个节点")
                return {"type": "graph_data", "data": graph}
            except Exception as e:
                logger.error(f"❌ [图检索] 失败: {e}")
                return {"type": "graph_data", "data": {"nodes": [], "edges": []}, "error": str(e)}

        # 并发执行检索任务
        search_tasks = [vector_search_task(), graph_search_task()]
        for coro in asyncio.as_completed(search_tasks):
            result = await coro
            yield json.dumps(result, ensure_ascii=False) + "\n"

        # 构建上下文
        await asyncio.sleep(0.5)
        context_parts = []

        if vector_chunks:
            context_parts.append("【相关文档片段】")
            for i, chunk in enumerate(vector_chunks, 1):
                distance = chunk.get('distance', None)
                distance_info = f" (相似度距离: {distance:.3f})" if distance is not None else ""
                context_parts.append(f"{i}. {chunk['content']}{distance_info}")

        if graph_data.get("nodes"):
            context_parts.append("\n【相关知识图谱】")
            for node in graph_data["nodes"][:10]:
                context_parts.append(f"- {node['label']} ({node['type']})")

        if not context_parts:
            prompt = f"""请回答以下问题。注意：未找到相关文档或知识图谱信息，请基于你的知识直接回答。

问题：{question}

请提供详细的答案："""
        else:
            context = "\n".join(context_parts)
            prompt = f"""请基于以下上下文回答问题。如果上下文中没有相关信息，请如实说明。

上下文：
{context}

问题：{question}

请提供详细的答案："""

        # 生成答案
        try:
            model_config = self.qa_clients.get(model_name or DEFAULT_QA_MODEL, {})
            use_qidian = model_config.get("use_qidian_api", False)
            use_yansd = model_config.get("use_yansd_api", False)

            stream_generator = None
            is_qidian_or_yansd = False

            if use_yansd:
                session_cookie = model_config.get("session_cookie", "")
                new_api_user = model_config.get("new_api_user", "")
                stream_generator = lambda: self.call_yansd_api_stream(prompt, qa_model, session_cookie, new_api_user)
                is_qidian_or_yansd = True
            elif use_qidian:
                satoken = model_config.get("satoken", "")
                model_id = model_config.get("model_id", 32)
                config_session_id = model_config.get("session_id", "")
                import uuid
                session_id = config_session_id or conversation_id or str(uuid.uuid4())
                stream_generator = lambda: self.call_qidian_api_stream(prompt, session_id, satoken, model_id)
                is_qidian_or_yansd = True
            else:
                messages = [{"role": "system", "content": "你是一个专业的知识问答助手。"}]
                if history:
                    recent_history = history[-10:] if len(history) > 10 else history
                    for msg in recent_history:
                        if msg.get("role") and msg.get("content"):
                            messages.append({"role": msg["role"], "content": msg["content"]})
                messages.append({"role": "user", "content": prompt})

                response = qa_client.chat.completions.create(
                    model=qa_model,
                    messages=messages,
                    temperature=0.3,
                    max_tokens=2000,
                    stream=True
                )
                stream_generator = lambda: response
                is_qidian_or_yansd = False

            # 流式输出
            if is_qidian_or_yansd:
                for text in stream_generator():
                    if text:
                        full_answer += text
                        yield json.dumps({"type": "answer", "content": text}, ensure_ascii=False) + "\n"
            else:
                for chunk in stream_generator():
                    if chunk.choices and len(chunk.choices) > 0 and chunk.choices[0].delta.content:
                        content = chunk.choices[0].delta.content
                        full_answer += content
                        yield json.dumps({"type": "answer", "content": content}, ensure_ascii=False) + "\n"

            yield json.dumps({"type": "answer_done"}, ensure_ascii=False) + "\n"

            # 保存对话
            if conversation_id:
                try:
                    save_conversation_message(conversation_id, "user", question)
                    save_conversation_message(conversation_id, "assistant", full_answer)
                except Exception as e:
                    logger.error(f"保存对话失败: {e}")

        except Exception as e:
            logger.error(f"答案生成失败: {e}")
            yield json.dumps({"type": "error", "error": str(e)}, ensure_ascii=False) + "\n"


# ==================== RRF 融合函数 ====================

def rrf_fusion(results_list: List[List[Dict[str, Any]]], k: int = 60) -> List[Dict[str, Any]]:
    """
    Reciprocal Rank Fusion (RRF) 融合多个排序结果
    
    Args:
        results_list: 多个检索结果列表的列表
        k: RRF 参数，默认 60
    
    Returns:
        融合后的排序结果列表
    """
    if not results_list:
        return []
    
    # 统计每个结果项的 RRF 分数
    score_map = {}  # content -> {rrf_score, original_result}
    
    for results in results_list:
        if not results:
            continue
        
        # 按排名计算 RRF 分数
        for rank, item in enumerate(results, 1):
            content = item.get("content", "")
            if not content:
                continue

            # 使用 doc_id + chunk_index 作为唯一标识
            metadata = item.get("metadata", {})
            unique_key = f"{metadata.get('doc_id', '')}_{metadata.get('chunk_index', '')}"

            rrf_score = 1.0 / (k + rank)

            if unique_key not in score_map:
                score_map[unique_key] = {
                    "rrf_score": 0,
                    "original_result": item
                }

            score_map[unique_key]["rrf_score"] += rrf_score
    
    # 按 RRF 分数排序
    sorted_results = sorted(
        score_map.values(),
        key=lambda x: x["rrf_score"],
        reverse=True
    )
    
    # 返回排序后的结果
    fusion_results = []
    for item in sorted_results:
        result = item["original_result"].copy()
        result["rrf_score"] = item["rrf_score"]
        result["fusion_method"] = "rrf"
        fusion_results.append(result)
    
    return fusion_results


# ==================== BM25 检索方法 ====================

def search_bm25_chunks(query: str, doc_id: Optional[str] = None, top_k: int = None) -> List[Dict[str, Any]]:
    """
    BM25 关键词检索
    
    Args:
        query: 搜索关键词
        doc_id: 可选的文档ID过滤
        top_k: 返回结果数量
    
    Returns:
        检索结果列表
    """
    if not BM25_ENABLED or not ES_AVAILABLE:
        return []
    
    if top_k is None:
        top_k = BM25_TOP_K
    
    try:
        results = search_kg_documents(query, doc_id=doc_id, size=top_k)

        # 格式化结果，添加 distance 字段以便与向量结果兼容
        formatted_results = []

        # 动态归一化：基于当前批次的最大分数
        max_score = max((r.get("bm25_score", 0) for r in results), default=0)

        for r in results:
            # 将 bm25_score 转换为伪 distance（分数越高越好，但 distance 越小越好）
            bm25_score = r.get("bm25_score", 0)
            normalized_score = bm25_score / max_score if max_score > 0 else 0
            distance = 1.0 - normalized_score

            formatted_results.append({
                "content": r.get("content", ""),
                "metadata": r.get("metadata", {}),
                "distance": distance,
                "bm25_score": bm25_score,
                "source": "bm25"
            })
        
        print(f"[BM25] 搜索 '{query[:30]}...' 找到 {len(formatted_results)} 条结果")
        return formatted_results
        
    except Exception as e:
        print(f"[BM25] 搜索失败: {e}")
        return []


# ==================== 混合检索方法 ====================

def hybrid_search_chunks(query: str, top_k: int = None) -> List[Dict[str, Any]]:
    """
    混合检索：向量检索 + BM25 关键词检索 + RRF 融合
    
    Args:
        query: 搜索关键词
        top_k: 返回结果数量
    
    Returns:
        融合后的检索结果列表
    """
    if top_k is None:
        top_k = VECTOR_SEARCH_TOP_K
    
    # 如果未启用 BM25，直接返回向量检索结果
    if not BM25_ENABLED or not ES_AVAILABLE:
        # 这会调用类的 search_similar_chunks 方法
        kg_svc = KnowledgeGraphService()
        return kg_svc.search_similar_chunks(query, top_k)
    
    print(f"[HYBRID] 开始混合检索: query='{query[:30]}...'")
    
    # 并行执行向量检索和 BM25 检索
    vector_results = []
    bm25_results = []
    
    try:
        # 向量检索（使用统一的 top_k）
        kg_svc = KnowledgeGraphService()
        vector_results = kg_svc.search_similar_chunks(query, top_k)
        print(f"[HYBRID] 向量检索找到 {len(vector_results)} 条结果")
    except Exception as e:
        print(f"[HYBRID] 向量检索失败: {e}")
    
    try:
        # BM25 检索（使用统一的 top_k）
        bm25_results = search_bm25_chunks(query, top_k=top_k)
        print(f"[HYBRID] BM25 检索找到 {len(bm25_results)} 条结果")
    except Exception as e:
        print(f"[HYBRID] BM25 检索失败: {e}")
    
    # RRF 融合
    if vector_results and bm25_results:
        # 合并两种结果
        all_results = [vector_results, bm25_results]
        fusion_results = rrf_fusion(all_results, k=RRF_K)
        print(f"[HYBRID] RRF 融合后共 {len(fusion_results)} 条结果")
        
        # 返回 top_k 结果
        return fusion_results[:top_k]
    elif vector_results:
        return vector_results[:top_k]
    elif bm25_results:
        return bm25_results[:top_k]
    else:
        return []


# ==================== 修改保存方法以同时支持 Milvus 和 ChromaDB ====================

# 保存到向量存储后同时索引到 ES（用于 BM25 检索）
async def _process_document_pipeline_async(
    self,
    file_path: str,
    doc_id: Optional[str] = None,
    include_triplets: bool = False,
    force_refresh: bool = False
) -> Dict[str, Any]:
    """统一文档处理入口，返回 text/chunks/triplets 减少重复计算"""
    if not hasattr(self, "pipeline_cache"):
        self.pipeline_cache = {}
    if not hasattr(self, "pipeline_cache_ttl"):
        self.pipeline_cache_ttl = 1800

    cached = None if force_refresh else self.pipeline_cache.get(doc_id)
    if cached and time.time() - cached.get("updated_at", 0) > self.pipeline_cache_ttl:
        self.pipeline_cache.pop(doc_id, None)
        cached = None

    has_required_stage = cached and (
        cached.get("triplets") is not None if include_triplets else cached.get("chunks") is not None
    )
    if has_required_stage:
        return {
            "text": cached.get("text", ""),
            "chunks": cached.get("chunks", []),
            "triplets": cached.get("triplets", [])
        }

    text = cached.get("text") if cached and cached.get("text") is not None else self.parse_document(file_path)
    chunks = cached.get("chunks") if cached and cached.get("chunks") is not None else self.split_text(text)
    triplets = cached.get("triplets") if cached and cached.get("triplets") is not None else None

    if include_triplets and triplets is None:
        triplets = await self.extract_batch_async(chunks)

    if doc_id:
        self.pipeline_cache[doc_id] = {
            "text": text,
            "chunks": chunks,
            "triplets": triplets,
            "updated_at": time.time()
        }

    return {
        "text": text,
        "chunks": chunks,
        "triplets": triplets or []
    }


def _save_triplets_to_neo4j_batched(self, triplets: List[Dict[str, Any]], doc_id: str):
    """按 label/relation 分组批量写入 Neo4j，减少 session.run 调用次数"""
    def sanitize_identifier(value: Optional[str], default: str) -> str:
        safe_value = "".join([c if c.isalnum() or c == "_" else "_" for c in (value or default)])
        if not safe_value or safe_value[0].isdigit():
            return default
        return safe_value

    normalized_triplets = []
    for triplet in triplets:
        head = (triplet.get("head") or "").strip()
        tail = (triplet.get("tail") or "").strip()
        if not head or not tail:
            continue

        normalized_triplets.append({
            "head": head,
            "tail": tail,
            "head_type": sanitize_identifier(triplet.get("head_type"), "Entity"),
            "tail_type": sanitize_identifier(triplet.get("tail_type"), "Entity"),
            "relation": sanitize_identifier(triplet.get("relation"), "RELATED_TO")
        })

    if not normalized_triplets:
        return

    head_groups = defaultdict(list)
    tail_groups = defaultdict(list)
    relation_groups = defaultdict(list)

    for row in normalized_triplets:
        head_groups[row["head_type"]].append({"head": row["head"]})
        tail_groups[row["tail_type"]].append({"tail": row["tail"]})
        relation_groups[(row["head_type"], row["tail_type"], row["relation"])].append({
            "head": row["head"],
            "tail": row["tail"]
        })

    with self.neo4j_driver.session() as session:
        session.run("MERGE (d:Document {id: $doc_id})", doc_id=doc_id)

        for head_type, rows in head_groups.items():
            session.run(
                f"""
                UNWIND $rows AS row
                MERGE (h:{head_type} {{name: row.head}})
                ON CREATE SET h.created_at = timestamp()
                WITH h
                MATCH (d:Document {{id: $doc_id}})
                MERGE (h)-[:FROM_DOCUMENT]->(d)
                """,
                rows=rows,
                doc_id=doc_id
            )

        for tail_type, rows in tail_groups.items():
            session.run(
                f"""
                UNWIND $rows AS row
                MERGE (t:{tail_type} {{name: row.tail}})
                ON CREATE SET t.created_at = timestamp()
                WITH t
                MATCH (d:Document {{id: $doc_id}})
                MERGE (t)-[:FROM_DOCUMENT]->(d)
                """,
                rows=rows,
                doc_id=doc_id
            )

        for (head_type, tail_type, relation), rows in relation_groups.items():
            session.run(
                f"""
                UNWIND $rows AS row
                MATCH (h:{head_type} {{name: row.head}})
                MATCH (t:{tail_type} {{name: row.tail}})
                MERGE (h)-[r:{relation}]->(t)
                ON CREATE SET r.created_at = timestamp()
                """,
                rows=rows
            )


_original_save_chunks_to_chromadb = KnowledgeGraphService.save_chunks_to_chromadb

def _save_chunks_with_es(self, chunks: List[Dict[str, Any]], doc_id: str):
    """
    将文本块保存到向量存储（自动选择 Milvus 或 ChromaDB，同步到 ES 用于 BM25）
    """
    # 根据配置选择使用 Milvus 或 ChromaDB
    if self.milvus_client and self.milvus_client.connected:
        # 使用 Milvus
        self.save_chunks_to_milvus(chunks, doc_id)
    else:
        # 使用 ChromaDB
        _original_save_chunks_to_chromadb(self, chunks, doc_id)
    
    # 同时索引到 Elasticsearch（用于 BM25 检索）
    if BM25_ENABLED and ES_AVAILABLE:
        try:
            index_kg_documents(chunks, doc_id)
            print(f"[BM25] 文档 {doc_id} 已同步到 Elasticsearch")
        except Exception as e:
            print(f"[BM25] 同步到 Elasticsearch 失败: {e}")

# 替换原方法
def _get_graph_data_partitioned(self, doc_id: Optional[str] = None, limit: int = 100) -> Dict[str, Any]:
    """Build graph data with document-based combo partitions for the frontend."""

    filter_clause = ""
    query_params: Dict[str, Any] = {"limit": limit}

    if doc_id:
        filter_clause = """
          AND (
            EXISTS { MATCH (n)-[:FROM_DOCUMENT]->(:Document {id: $doc_id}) }
            OR EXISTS { MATCH (m)-[:FROM_DOCUMENT]->(:Document {id: $doc_id}) }
          )
        """
        query_params["doc_id"] = doc_id

    query = f"""
        MATCH (n)-[r]->(m)
        WHERE type(r) <> 'FROM_DOCUMENT'
          AND NOT 'Document' IN labels(n)
          AND NOT 'Document' IN labels(m)
          {filter_clause}
        OPTIONAL MATCH (n)-[:FROM_DOCUMENT]->(dn:Document)
        OPTIONAL MATCH (m)-[:FROM_DOCUMENT]->(dm:Document)
        WITH n, r, m,
             [doc_ref IN collect(DISTINCT dn.id) WHERE doc_ref IS NOT NULL] AS n_doc_ids,
             [doc_ref IN collect(DISTINCT dm.id) WHERE doc_ref IS NOT NULL] AS m_doc_ids
        RETURN n, r, m, n_doc_ids, m_doc_ids
        LIMIT $limit
    """

    def normalize_doc_ids(raw_doc_ids: Optional[List[Any]]) -> List[str]:
        doc_refs: List[str] = []
        for raw_doc_id in raw_doc_ids or []:
            doc_ref = str(raw_doc_id).strip()
            if not doc_ref or doc_ref in doc_refs:
                continue
            doc_refs.append(doc_ref)

        if doc_id and doc_id in doc_refs:
            doc_refs = [doc_id] + [ref for ref in doc_refs if ref != doc_id]

        return doc_refs

    def build_combo_id(doc_ref: Optional[str]) -> Optional[str]:
        return f"doc::{doc_ref}" if doc_ref else None

    def build_combo_label(doc_ref: Optional[str]) -> str:
        if not doc_ref:
            return "Ungrouped"
        prefix = "Current Doc" if doc_id and doc_ref == doc_id else "Doc"
        return f"{prefix} {doc_ref[:8]}"

    with self.neo4j_driver.session() as session:
        result = session.run(query, **query_params)

        nodes: Dict[str, Dict[str, Any]] = {}
        edges: List[Dict[str, Any]] = []
        edge_keys = set()

        def ensure_partitioned_node(entity, partition_doc_id: Optional[str], source_doc_ids: List[str]) -> str:
            node_id = f"{entity.element_id}::{partition_doc_id}" if partition_doc_id else entity.element_id
            if node_id in nodes:
                return node_id

            combo_id = build_combo_id(partition_doc_id)
            combo_label = build_combo_label(partition_doc_id)
            nodes[node_id] = {
                "id": node_id,
                "entity_id": entity.element_id,
                "label": entity.get("name", "Unknown"),
                "type": list(entity.labels)[0] if entity.labels else "Entity",
                "combo_id": combo_id,
                "combo_label": combo_label if combo_id else None,
                "source_document_ids": source_doc_ids,
                "document_count": len(source_doc_ids)
            }
            return node_id

        def add_edge(source_id: str, target_id: str, relation_label: str, is_cross_combo: bool):
            edge_key = (source_id, target_id, relation_label, is_cross_combo)
            if edge_key in edge_keys:
                return

            edge_keys.add(edge_key)
            edges.append({
                "id": f"{source_id}->{relation_label}->{target_id}",
                "source": source_id,
                "target": target_id,
                "label": relation_label,
                "is_cross_combo": is_cross_combo
            })

        for record in result:
            n = record["n"]
            r = record["r"]
            m = record["m"]

            n_doc_ids = normalize_doc_ids(record.get("n_doc_ids"))
            m_doc_ids = normalize_doc_ids(record.get("m_doc_ids"))

            if doc_id and not n_doc_ids:
                n_doc_ids = [doc_id]
            if doc_id and not m_doc_ids:
                m_doc_ids = [doc_id]

            shared_doc_ids = [doc_ref for doc_ref in n_doc_ids if doc_ref in m_doc_ids]

            if shared_doc_ids:
                for shared_doc_id in shared_doc_ids:
                    source_id = ensure_partitioned_node(n, shared_doc_id, n_doc_ids)
                    target_id = ensure_partitioned_node(m, shared_doc_id, m_doc_ids)
                    add_edge(source_id, target_id, r.type, False)
                continue

            source_partition = n_doc_ids[0] if n_doc_ids else None
            target_partition = m_doc_ids[0] if m_doc_ids else None
            source_id = ensure_partitioned_node(n, source_partition, n_doc_ids)
            target_id = ensure_partitioned_node(m, target_partition, m_doc_ids)
            add_edge(source_id, target_id, r.type, source_partition != target_partition)

        return {"nodes": list(nodes.values()), "edges": edges}


def _apply_graph_edits(
    self,
    nodes: List[Dict[str, Any]],
    edges: List[Dict[str, Any]]
) -> Dict[str, Any]:
    """Apply question-time graph edits back into Neo4j."""

    def sanitize_identifier(value: Optional[str], default: str) -> str:
        safe_value = "".join([c if c.isalnum() or c == "_" else "_" for c in (value or default)])
        if not safe_value or safe_value[0].isdigit():
            return default
        return safe_value

    normalized_nodes: List[Dict[str, str]] = []
    for node in nodes or []:
        node_id = str(node.get("id") or "").strip()
        name = str(node.get("label") or node.get("name") or "").strip()
        node_type = sanitize_identifier(node.get("type"), "Entity")
        current_type = sanitize_identifier(node.get("current_type") or node.get("type"), "Entity")
        if not node_id or not name:
            continue
        normalized_nodes.append({
            "id": node_id,
            "name": name,
            "type": node_type,
            "current_type": current_type
        })

    normalized_edges: List[Dict[str, str]] = []
    for edge in edges or []:
        edge_id = str(edge.get("id") or "").strip()
        source = str(edge.get("source") or "").strip()
        target = str(edge.get("target") or "").strip()
        relation = sanitize_identifier(edge.get("label"), "RELATED_TO")
        current_relation = sanitize_identifier(edge.get("current_label") or edge.get("label"), "RELATED_TO")
        if not edge_id or not source or not target:
            continue
        normalized_edges.append({
            "id": edge_id,
            "source": source,
            "target": target,
            "relation": relation,
            "current_relation": current_relation
        })

    updated_nodes = 0
    updated_edges = 0

    with self.neo4j_driver.session() as session:
        for node in normalized_nodes:
            session.run(
                """
                MATCH (n)
                WHERE elementId(n) = $node_id
                  AND NOT 'Document' IN labels(n)
                SET n.name = $name,
                    n.updated_at = timestamp()
                """,
                node_id=node["id"],
                name=node["name"]
            )
            updated_nodes += 1

            if node["type"] != node["current_type"]:
                session.run(
                    f"""
                    MATCH (n:{node['current_type']})
                    WHERE elementId(n) = $node_id
                    REMOVE n:{node['current_type']}
                    SET n:{node['type']}
                    """,
                    node_id=node["id"]
                )

        for edge in normalized_edges:
            session.run(
                """
                MATCH ()-[r]->()
                WHERE elementId(r) = $edge_id
                DELETE r
                """,
                edge_id=edge["id"]
            )
            session.run(
                f"""
                MATCH (s)
                WHERE elementId(s) = $source_id
                MATCH (t)
                WHERE elementId(t) = $target_id
                MERGE (s)-[r:{edge['relation']}]->(t)
                SET r.updated_at = timestamp()
                """,
                source_id=edge["source"],
                target_id=edge["target"]
            )
            updated_edges += 1

    return {
        "updated_nodes": updated_nodes,
        "updated_edges": updated_edges
    }


KnowledgeGraphService.process_document_pipeline_async = _process_document_pipeline_async
KnowledgeGraphService.save_triplets_to_neo4j = _save_triplets_to_neo4j_batched
KnowledgeGraphService.save_chunks_to_chromadb = _save_chunks_with_es
KnowledgeGraphService.get_graph_data = _get_graph_data_partitioned
KnowledgeGraphService.apply_graph_edits = _apply_graph_edits


# 全局服务实例
kg_service = KnowledgeGraphService()
